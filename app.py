
from flask import Flask, request, jsonify
from flask_cors import CORS
from docx import Document
from datetime import datetime
import os
import openai
import tempfile
import pdfplumber
import docx2txt
import pandas as pd
import pptx
from PIL import Image
import pytesseract
import io

app = Flask(__name__)
CORS(app)

# Ensure the reports directory exists
REPORT_FOLDER = os.path.join(app.root_path, 'static', 'reports')
os.makedirs(REPORT_FOLDER, exist_ok=True)

# Initialize OpenAI
openai.api_key = os.getenv("OPENAI_API_KEY")

def extract_text(file_storage):
    filename = file_storage.filename.lower()
    if filename.endswith(".pdf"):
        with pdfplumber.open(file_storage.stream) as pdf:
            return "\n".join(page.extract_text() or "" for page in pdf.pages)
    elif filename.endswith(".docx"):
        with tempfile.NamedTemporaryFile(delete=False, suffix=".docx") as tmp:
            file_storage.save(tmp.name)
            return docx2txt.process(tmp.name)
    elif filename.endswith(".pptx"):
        presentation = pptx.Presentation(file_storage)
        text = []
        for slide in presentation.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text.append(shape.text)
        return "\n".join(text)
    elif filename.endswith(".png") or filename.endswith(".jpg") or filename.endswith(".jpeg"):
        image = Image.open(file_storage.stream)
        return pytesseract.image_to_string(image)
    elif filename.endswith(".xlsx"):
        xls = pd.read_excel(file_storage, sheet_name=None)
        return "\n".join([df.to_string() for df in xls.values()])
    elif filename.endswith(".csv"):
        df = pd.read_csv(file_storage)
        return df.to_string()
    else:
        return ""

def generate_section(title, instruction, context):
    try:
        prompt = f"{instruction}\n\nContext:\n{context}"
        response = openai.chat.completions.create(
            model="gpt-4",
            messages=[{"role": "user", "content": prompt}],
            temperature=0.7
        )
        return response.choices[0].message.content.strip()
    except Exception as e:
        return f"Error generating this section: {e}"

@app.route('/')
def home():
    return "Management Optimization Backend is Live!"

@app.route('/generate', methods=['POST'])
def generate_report():
    files = [request.files.get('file1'), request.files.get('file2'), request.files.get('file3')]
    context = ""

    for file in files:
        if file:
            context += extract_text(file) + "\n"

    if not context.strip():
        return jsonify({"error": "No valid file content found."}), 400

    doc = Document()
    doc.add_heading("Management Optimization Report", 0)

    sections = [
        ("Executive Summary", "Summarize the current management practices, leadership style, and team dynamics based on the content."),
        ("1. Leadership & Team Structure", "Evaluate the clarity of leadership roles, the hierarchy, and how effectively the team operates within the structure."),
        ("2. Workflow & Productivity", "Analyze internal workflows and productivity levels, identifying inefficiencies or bottlenecks."),
        ("3. Decision-Making Processes", "Assess how decisions are made within the organization and evaluate their impact on overall management effectiveness."),
        ("4. Management KPIs & Performance", "Identify key performance indicators used to measure managerial success and analyze whether they align with business goals."),
        ("5. Recommendations & Optimization", "Based on the findings, provide tailored recommendations for improving management efficiency and team output."),
        ("Conclusion", "Conclude with a high-level evaluation of the management landscape and recommended next steps.")
    ]

    for title, instruction in sections:
        doc.add_heading(title, level=1)
        generated = generate_section(title, instruction, context)
        doc.add_paragraph(generated)

    filename = f"management_report_{datetime.now().strftime('%Y%m%d%H%M%S')}.docx"
    file_path = os.path.join(REPORT_FOLDER, filename)
    doc.save(file_path)

    return jsonify({'download_url': f'/static/reports/{filename}'})

if __name__ == '__main__':
    port = int(os.environ.get("PORT", 5000))
    app.run(host='0.0.0.0', port=port)
